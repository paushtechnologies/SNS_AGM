
import collections
import collections.abc
from pptx import Presentation
import json
import os

def extract_pptx_content(pptx_path):
    prs = Presentation(pptx_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "title": "",
            "content": []
        }
        
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    slide_data["content"].append(shape.text.strip())
        
        slides_content.append(slide_data)
    
    return slides_content

if __name__ == "__main__":
    pptx_file = r"C:\Users\asus\Downloads\SNS\ABPM 2026-27 Draft 2.pptx"
    content = extract_pptx_content(pptx_file)
    with open("slides_data.json", "w", encoding="utf-8") as f:
        json.dump(content, f, indent=2)
