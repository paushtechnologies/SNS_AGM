import sys
from pptx import Presentation

def extract_pptx_content(pptx_file):
    prs = Presentation(pptx_file)
    content = ""
    for i, slide in enumerate(prs.slides):
        content += f"--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
        content += "\n"
    return content

if __name__ == "__main__":
    pptx_path = r"C:\Users\asus\Downloads\SNS\SNS_10000_Crore_Mindset - Nasik .pptx"
    output_path = r"c:\Users\asus\Downloads\SNS\Website\elevate_scale_content.txt"
    try:
        extracted_text = extract_pptx_content(pptx_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)
        print(f"Content extracted to {output_path}")
    except Exception as e:
        print(f"Error: {e}")
