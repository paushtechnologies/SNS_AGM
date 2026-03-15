
import collections
import collections.abc
from pptx import Presentation
import json

def extract_pptx_content(pptx_path):
    prs = Presentation(pptx_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "title": "",
            "content": []
        }
        
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    slide_data["content"].append(shape.text.strip())
        
        slides_content.append(slide_data)
    
    return slides_content

if __name__ == "__main__":
    pptx_file = r"C:\Users\asus\Downloads\SNS\SNS_10000_Crore_Mindset - Nasik .pptx"
    content = extract_pptx_content(pptx_file)
    with open("session04_data.json", "w", encoding="utf-8") as f:
        json.dump(content, f, indent=2)
    print("Content saved to session04_data.json")
